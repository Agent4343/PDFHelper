"""Comprehensive integration tests for the Procedure Builder.

Tests cover: presets, CRUD, filters, status workflow, file attachments
(DOCX, PDF, PPTX, TXT, CSV, rejection of unsupported types), logo upload,
preview, route ordering, and frontend HTML verification.
"""

import io
import pytest


# ---------------------------------------------------------------------------
# Helpers — minimal valid test files
# ---------------------------------------------------------------------------

_MINIMAL_PDF = (
    b"%PDF-1.4\n"
    b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
    b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
    b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R>>endobj\n"
    b"xref\n0 4\n"
    b"0000000000 65535 f \n"
    b"0000000009 00000 n \n"
    b"0000000058 00000 n \n"
    b"0000000115 00000 n \n"
    b"trailer<</Size 4/Root 1 0 R>>\nstartxref\n190\n%%EOF"
)


def _make_docx_bytes():
    from docx import Document
    doc = Document()
    doc.add_heading("Test Procedure", level=1)
    doc.add_paragraph("Step 1: OPEN valve.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Safety Briefing"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_small_png():
    """1x1 red PNG."""
    import struct, zlib
    raw = b"\x00\xff\x00\x00"
    compressed = zlib.compress(raw)
    def chunk(ctype, data):
        c = ctype + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", compressed)
        + chunk(b"IEND", b"")
    )


# ---------------------------------------------------------------------------
# Fixture: create a procedure session and return its ID
# ---------------------------------------------------------------------------

@pytest.fixture
async def proc_session(client, api_headers):
    """Create a procedure session and return its ID."""
    r = await client.post(
        "/procedures",
        headers=api_headers,
        json={"title": "Test Procedure", "facility": "Hebron", "craft": "Operations", "category": "Maintenance"},
    )
    assert r.status_code == 200
    return r.json()["id"]


# ===========================================================================
# 1. Facility presets
# ===========================================================================

@pytest.mark.anyio
async def test_get_presets(client, api_headers):
    r = await client.get("/procedures/presets", headers=api_headers)
    assert r.status_code == 200
    data = r.json()
    assert "Hebron" in data["facilities"]
    assert "Hibernia" in data["facilities"]
    assert "Operations" in data["facilities"]["Hebron"]["crafts"]
    assert data["facilities"]["Hebron"]["doc_prefix"] == "CAHE-EC-OOPRO"


# ===========================================================================
# 2. CRUD — Create, list, get, delete
# ===========================================================================

@pytest.mark.anyio
async def test_create_procedure_new(client, api_headers):
    r = await client.post(
        "/procedures",
        headers=api_headers,
        json={"title": "Valve Lineup", "facility": "Hebron", "craft": "Operations"},
    )
    assert r.status_code == 200
    data = r.json()
    assert data["title"] == "Valve Lineup"
    assert data["status"] == "gathering"
    assert data["mode"] == "new"
    assert len(data["messages"]) == 1
    assert data["messages"][0]["role"] == "assistant"
    assert "Hebron" in data["messages"][0]["content"]


@pytest.mark.anyio
async def test_create_procedure_no_facility(client, api_headers):
    r = await client.post(
        "/procedures",
        headers=api_headers,
        json={"title": "Generic Proc"},
    )
    assert r.status_code == 200
    data = r.json()
    assert "facility" in data["messages"][0]["content"].lower() or "Let's start" in data["messages"][0]["content"]


@pytest.mark.anyio
async def test_create_procedure_default_title(client, api_headers):
    r = await client.post("/procedures", headers=api_headers, json={})
    assert r.status_code == 200
    assert r.json()["title"] == "New Procedure"


@pytest.mark.anyio
async def test_list_procedures(client, api_headers, proc_session):
    r = await client.get("/procedures", headers=api_headers)
    assert r.status_code == 200
    data = r.json()
    assert "sessions" in data
    assert "total" in data
    assert data["total"] >= 1


@pytest.mark.anyio
async def test_list_procedures_pagination(client, api_headers, proc_session):
    r = await client.get("/procedures?skip=0&limit=5", headers=api_headers)
    assert r.status_code == 200
    assert len(r.json()["sessions"]) <= 5


@pytest.mark.anyio
async def test_get_procedure(client, api_headers, proc_session):
    r = await client.get(f"/procedures/{proc_session}", headers=api_headers)
    assert r.status_code == 200
    data = r.json()
    assert data["id"] == proc_session
    assert data["status"] == "gathering"
    assert data["facility"] == "Hebron"
    assert data["category"] == "Maintenance"
    assert len(data["messages"]) >= 1


@pytest.mark.anyio
async def test_get_procedure_not_found(client, api_headers):
    r = await client.get("/procedures/nonexistent-id", headers=api_headers)
    assert r.status_code == 404


@pytest.mark.anyio
async def test_delete_procedure(client, api_headers):
    r = await client.post("/procedures", headers=api_headers, json={"title": "Delete Me"})
    sid = r.json()["id"]
    r = await client.delete(f"/procedures/{sid}", headers=api_headers)
    assert r.status_code == 200
    assert "deleted" in r.json()["detail"].lower()
    r = await client.get(f"/procedures/{sid}", headers=api_headers)
    assert r.status_code == 404


@pytest.mark.anyio
async def test_delete_procedure_not_found(client, api_headers):
    r = await client.delete("/procedures/nonexistent-id", headers=api_headers)
    assert r.status_code == 404


# ===========================================================================
# 3. Filters — facility, category, status
# ===========================================================================

@pytest.mark.anyio
async def test_filter_by_facility(client, api_headers, proc_session):
    r = await client.get("/procedures?facility=Hebron", headers=api_headers)
    assert r.status_code == 200
    for s in r.json()["sessions"]:
        assert s["facility"] == "Hebron"


@pytest.mark.anyio
async def test_filter_by_category(client, api_headers, proc_session):
    r = await client.get("/procedures?category=Maintenance", headers=api_headers)
    assert r.status_code == 200
    for s in r.json()["sessions"]:
        assert s["category"] == "Maintenance"


@pytest.mark.anyio
async def test_filter_by_status(client, api_headers, proc_session):
    r = await client.get("/procedures?status=gathering", headers=api_headers)
    assert r.status_code == 200
    for s in r.json()["sessions"]:
        assert s["status"] == "gathering"


@pytest.mark.anyio
async def test_filter_no_match(client, api_headers):
    r = await client.get("/procedures?facility=NonexistentPlatform", headers=api_headers)
    assert r.status_code == 200
    assert r.json()["total"] == 0


# ===========================================================================
# 4. Status workflow transitions
# ===========================================================================

@pytest.mark.anyio
async def test_status_gathering_to_drafting(client, api_headers, proc_session):
    r = await client.post(
        f"/procedures/{proc_session}/status",
        headers=api_headers,
        json={"status": "drafting"},
    )
    assert r.status_code == 200
    assert r.json()["status"] == "drafting"


@pytest.mark.anyio
async def test_status_gathering_to_complete(client, api_headers, proc_session):
    r = await client.post(
        f"/procedures/{proc_session}/status",
        headers=api_headers,
        json={"status": "complete"},
    )
    assert r.status_code == 200
    assert r.json()["status"] == "complete"


@pytest.mark.anyio
async def test_status_full_workflow(client, api_headers):
    """Test full lifecycle: gathering -> drafting -> review -> approved -> complete."""
    r = await client.post("/procedures", headers=api_headers, json={"title": "Lifecycle Test"})
    sid = r.json()["id"]

    transitions = [
        ("drafting", 200),
        ("review", 200),
        ("approved", 200),
        ("complete", 200),
    ]
    for status, expected_code in transitions:
        r = await client.post(
            f"/procedures/{sid}/status",
            headers=api_headers,
            json={"status": status},
        )
        assert r.status_code == expected_code, f"Failed transition to {status}: {r.text}"
        assert r.json()["status"] == status


@pytest.mark.anyio
async def test_status_invalid_transition(client, api_headers, proc_session):
    r = await client.post(
        f"/procedures/{proc_session}/status",
        headers=api_headers,
        json={"status": "approved"},
    )
    assert r.status_code == 400
    assert "Cannot transition" in r.json()["detail"]


@pytest.mark.anyio
async def test_status_review_cannot_go_to_complete(client, api_headers):
    """review -> complete is not a valid transition."""
    r = await client.post("/procedures", headers=api_headers, json={"title": "Review Block Test"})
    sid = r.json()["id"]
    await client.post(f"/procedures/{sid}/status", headers=api_headers, json={"status": "drafting"})
    await client.post(f"/procedures/{sid}/status", headers=api_headers, json={"status": "review"})
    r = await client.post(f"/procedures/{sid}/status", headers=api_headers, json={"status": "complete"})
    assert r.status_code == 400


@pytest.mark.anyio
async def test_status_complete_to_review(client, api_headers):
    """complete -> review is valid (reopen for re-review)."""
    r = await client.post("/procedures", headers=api_headers, json={"title": "Reopen Test"})
    sid = r.json()["id"]
    await client.post(f"/procedures/{sid}/status", headers=api_headers, json={"status": "complete"})
    r = await client.post(f"/procedures/{sid}/status", headers=api_headers, json={"status": "review"})
    assert r.status_code == 200
    assert r.json()["status"] == "review"


@pytest.mark.anyio
async def test_status_not_found(client, api_headers):
    r = await client.post(
        "/procedures/nonexistent-id/status",
        headers=api_headers,
        json={"status": "drafting"},
    )
    assert r.status_code == 404


# ===========================================================================
# 5. File attachments
# ===========================================================================

@pytest.mark.anyio
async def test_attach_txt(client, api_headers, proc_session):
    files = [("file", ("notes.txt", io.BytesIO(b"Safety valve check procedure"), "text/plain"))]
    r = await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 200
    data = r.json()
    assert data["filename"] == "notes.txt"
    assert data["chars"] > 0


@pytest.mark.anyio
async def test_attach_csv(client, api_headers, proc_session):
    csv_data = b"Step,Action,Who\n1,OPEN valve,Ops Tech\n2,VERIFY flow,Ops Tech"
    files = [("file", ("steps.csv", io.BytesIO(csv_data), "text/csv"))]
    r = await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 200
    assert r.json()["filename"] == "steps.csv"


@pytest.mark.anyio
async def test_attach_pdf(client, api_headers, proc_session):
    files = [("file", ("ref.pdf", io.BytesIO(_MINIMAL_PDF), "application/pdf"))]
    r = await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 200
    assert r.json()["filename"] == "ref.pdf"


@pytest.mark.anyio
async def test_attach_docx(client, api_headers, proc_session):
    docx_bytes = _make_docx_bytes()
    files = [("file", ("procedure.docx", io.BytesIO(docx_bytes), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))]
    r = await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 200
    data = r.json()
    assert data["filename"] == "procedure.docx"
    assert data["chars"] > 0


@pytest.mark.anyio
async def test_attach_pptx(client, api_headers, proc_session):
    """PPTX files should be accepted and text extracted."""
    pptx_bytes = _make_pptx_bytes()
    files = [("file", ("slides.pptx", io.BytesIO(pptx_bytes), "application/vnd.openxmlformats-officedocument.presentationml.presentation"))]
    r = await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 200
    data = r.json()
    assert data["filename"] == "slides.pptx"
    assert data["chars"] > 0


@pytest.mark.anyio
async def test_attach_pptx_text_extraction(client, api_headers, proc_session):
    """Verify PPTX text content is extracted and stored in the message."""
    pptx_bytes = _make_pptx_bytes()
    files = [("file", ("brief.pptx", io.BytesIO(pptx_bytes), "application/vnd.openxmlformats-officedocument.presentationml.presentation"))]
    await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    r = await client.get(f"/procedures/{proc_session}", headers=api_headers)
    messages = r.json()["messages"]
    attach_msgs = [m for m in messages if m["role"] == "user" and "Attached file" in m["content"]]
    assert len(attach_msgs) >= 1
    assert "Safety Briefing" in attach_msgs[-1]["content"]


@pytest.mark.anyio
async def test_attach_unsupported_format(client, api_headers, proc_session):
    files = [("file", ("image.jpg", io.BytesIO(b"\xff\xd8\xff\xe0"), "image/jpeg"))]
    r = await client.post(
        f"/procedures/{proc_session}/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 400
    assert ".pptx" in r.json()["detail"]


@pytest.mark.anyio
async def test_attach_not_found_session(client, api_headers):
    files = [("file", ("notes.txt", io.BytesIO(b"hello"), "text/plain"))]
    r = await client.post(
        "/procedures/nonexistent-id/attach",
        headers=api_headers,
        files=files,
    )
    assert r.status_code == 404


# ===========================================================================
# 6. Logo upload
# ===========================================================================

@pytest.mark.anyio
async def test_upload_logo_png(client, api_headers):
    png = _make_small_png()
    files = [("file", ("logo.png", io.BytesIO(png), "image/png"))]
    r = await client.post("/procedures/logo", headers=api_headers, files=files)
    assert r.status_code == 200
    assert "Logo uploaded" in r.json()["detail"]


@pytest.mark.anyio
async def test_get_logo_status(client, api_headers):
    r = await client.get("/procedures/logo", headers=api_headers)
    assert r.status_code == 200
    assert "has_logo" in r.json()


@pytest.mark.anyio
async def test_upload_logo_invalid_format(client, api_headers):
    files = [("file", ("logo.gif", io.BytesIO(b"GIF89a"), "image/gif"))]
    r = await client.post("/procedures/logo", headers=api_headers, files=files)
    assert r.status_code == 400
    assert "PNG/JPG" in r.json()["detail"]


@pytest.mark.anyio
async def test_upload_logo_too_large(client, api_headers):
    big = b"\x00" * (2 * 1024 * 1024 + 1)
    files = [("file", ("logo.png", io.BytesIO(big), "image/png"))]
    r = await client.post("/procedures/logo", headers=api_headers, files=files)
    assert r.status_code == 400
    assert "too large" in r.json()["detail"].lower()


# ===========================================================================
# 7. Route ordering — /procedures/logo must not be shadowed by {session_id}
# ===========================================================================

@pytest.mark.anyio
async def test_logo_route_not_shadowed(client, api_headers):
    """GET /procedures/logo must return logo status, not 404 from {session_id}."""
    r = await client.get("/procedures/logo", headers=api_headers)
    assert r.status_code == 200
    assert "has_logo" in r.json()


@pytest.mark.anyio
async def test_presets_route_not_shadowed(client, api_headers):
    """GET /procedures/presets must return facility data, not 404 from {session_id}."""
    r = await client.get("/procedures/presets", headers=api_headers)
    assert r.status_code == 200
    assert "facilities" in r.json()


# ===========================================================================
# 8. Upload .docx for procedure update mode
# ===========================================================================

@pytest.mark.anyio
async def test_upload_procedure_docx(client, api_headers):
    docx_bytes = _make_docx_bytes()
    files = [("file", ("valve-proc.docx", io.BytesIO(docx_bytes), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))]
    r = await client.post("/procedures/upload", headers=api_headers, files=files)
    assert r.status_code == 200
    data = r.json()
    assert data["mode"] == "update"
    assert data["status"] == "gathering"
    assert "Test Procedure" in data["title"] or "valve-proc" in data["title"]


@pytest.mark.anyio
async def test_upload_procedure_non_docx_rejected(client, api_headers):
    files = [("file", ("notes.txt", io.BytesIO(b"hello"), "text/plain"))]
    r = await client.post("/procedures/upload", headers=api_headers, files=files)
    assert r.status_code == 400
    assert ".docx" in r.json()["detail"]


# ===========================================================================
# 9. Preview — requires generated output
# ===========================================================================

@pytest.mark.anyio
async def test_preview_no_output(client, api_headers, proc_session):
    """Preview should fail if procedure hasn't been generated yet."""
    r = await client.get(f"/procedures/{proc_session}/preview", headers=api_headers)
    assert r.status_code == 400
    assert "not yet generated" in r.json()["detail"].lower()


@pytest.mark.anyio
async def test_preview_not_found(client, api_headers):
    r = await client.get("/procedures/nonexistent-id/preview", headers=api_headers)
    assert r.status_code == 404


# ===========================================================================
# 10. Chat — requires auth, validates message
# ===========================================================================

@pytest.mark.anyio
async def test_chat_requires_auth(client, proc_session):
    """POST /procedures/{id}/chat without auth returns 401."""
    r = await client.post(
        f"/procedures/{proc_session}/chat",
        json={"message": "hello"},
    )
    assert r.status_code == 401


@pytest.mark.anyio
async def test_chat_empty_message(client, api_headers, proc_session):
    r = await client.post(
        f"/procedures/{proc_session}/chat",
        headers=api_headers,
        json={"message": ""},
    )
    assert r.status_code == 400
    assert "required" in r.json()["detail"].lower()


@pytest.mark.anyio
async def test_chat_not_found(client, api_headers):
    r = await client.post(
        "/procedures/nonexistent-id/chat",
        headers=api_headers,
        json={"message": "hello"},
    )
    assert r.status_code == 404


# ===========================================================================
# 11. Compliance — requires content
# ===========================================================================

@pytest.mark.anyio
async def test_compliance_no_content(client, api_headers, proc_session):
    r = await client.post(f"/procedures/{proc_session}/compliance", headers=api_headers)
    assert r.status_code == 400
    assert "No procedure content" in r.json()["detail"]


@pytest.mark.anyio
async def test_compliance_not_found(client, api_headers):
    r = await client.post("/procedures/nonexistent-id/compliance", headers=api_headers)
    assert r.status_code == 404


# ===========================================================================
# 12. Auth gating — all procedure endpoints need credentials
# ===========================================================================

@pytest.mark.anyio
async def test_procedures_list_requires_auth(client):
    r = await client.get("/procedures")
    assert r.status_code == 401


@pytest.mark.anyio
async def test_procedures_create_requires_auth(client):
    r = await client.post("/procedures", json={"title": "No auth"})
    assert r.status_code == 401


@pytest.mark.anyio
async def test_procedures_delete_requires_auth(client):
    r = await client.delete("/procedures/some-id")
    assert r.status_code == 401


@pytest.mark.anyio
async def test_procedures_attach_requires_auth(client):
    files = [("file", ("n.txt", io.BytesIO(b"x"), "text/plain"))]
    r = await client.post("/procedures/some-id/attach", files=files)
    assert r.status_code == 401


# ===========================================================================
# 13. Frontend HTML checks
# ===========================================================================

@pytest.mark.anyio
async def test_frontend_pptx_accept(client):
    """The file-attach input should accept .pptx files."""
    r = await client.get("/")
    assert r.status_code == 200
    html = r.text
    assert ".pptx" in html


@pytest.mark.anyio
async def test_frontend_has_facility_dropdown(client):
    r = await client.get("/")
    html = r.text
    assert "proc-facility" in html


@pytest.mark.anyio
async def test_frontend_has_status_controls(client):
    r = await client.get("/")
    html = r.text
    assert "changeProcStatus" in html or "proc-status" in html


@pytest.mark.anyio
async def test_frontend_has_logo_upload(client):
    r = await client.get("/")
    html = r.text
    assert "uploadLogo" in html


@pytest.mark.anyio
async def test_frontend_has_compliance_button(client):
    r = await client.get("/")
    html = r.text
    assert "compliance" in html.lower()
