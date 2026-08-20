import hashlib
import io
import json
import os
import re
import uuid
from datetime import datetime, timezone, timedelta
from pathlib import Path

from fastapi import HTTPException, UploadFile

from config import ENCRYPTION_KEY, UPLOAD_DIR, MAX_FILE_SIZE, AUTO_CLEANUP_HOURS, CHAT_MODEL, AGENT_MAX_TOKENS, AGENT_MODELS
from database import SessionLocal, DBDocument, DBAgentCache
from ocr import extract_text_with_ocr_fallback, extract_structured_text

# Re-export for backward compat
from utils import is_retryable_error, parse_json_response

# ---------------------------------------------------------------------------
# File type constants
# ---------------------------------------------------------------------------

PDF_MAGIC_BYTES = b"%PDF-"
ALLOWED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp"}
ALLOWED_SPREADSHEET_EXTENSIONS = {".xlsx", ".xls", ".csv"}
ALLOWED_WORD_EXTENSIONS = {".docx", ".doc"}
ALLOWED_PRESENTATION_EXTENSIONS = {".pptx"}
ALLOWED_TEXT_EXTENSIONS = {
    ".js", ".html", ".htm", ".css", ".md", ".txt", ".json", ".xml",
    ".yaml", ".yml", ".py", ".ts", ".tsx", ".jsx", ".sql", ".sh",
    ".env", ".toml", ".ini", ".cfg",
}


# ---------------------------------------------------------------------------
# File helpers
# ---------------------------------------------------------------------------

def _sanitize_filename(filename: str) -> str:
    name = Path(filename).name
    name = re.sub(r"[^\w\-. ]", "_", name)
    name = name.lstrip(".")
    return name or "unnamed.pdf"


def _validate_filepath(filepath: Path) -> Path:
    try:
        resolved = filepath.resolve()
        if not resolved.is_relative_to(UPLOAD_DIR.resolve()):
            raise ValueError(f"Path escapes upload directory: {filepath}")
        if resolved.is_symlink():
            raise ValueError(f"Symlinks not allowed: {filepath}")
        return resolved
    except (OSError, ValueError):
        raise


def _safe_unlink(filepath: Path) -> None:
    try:
        validated = _validate_filepath(filepath)
        validated.unlink()
    except FileNotFoundError:
        pass
    except ValueError:
        import logging
        logging.getLogger("pdfhelper").warning("Blocked path traversal attempt: %s", filepath)


def _verify_pdf_content(data: bytes) -> bool:
    return data[:5] == PDF_MAGIC_BYTES


def _is_image_file(filename: str) -> bool:
    return any(filename.lower().endswith(ext) for ext in ALLOWED_IMAGE_EXTENSIONS)


def _is_spreadsheet_file(filename: str) -> bool:
    return any(filename.lower().endswith(ext) for ext in ALLOWED_SPREADSHEET_EXTENSIONS)


def _is_word_file(filename: str) -> bool:
    return any(filename.lower().endswith(ext) for ext in ALLOWED_WORD_EXTENSIONS)


def _is_presentation_file(filename: str) -> bool:
    return any(filename.lower().endswith(ext) for ext in ALLOWED_PRESENTATION_EXTENSIONS)


def _is_text_file(filename: str) -> bool:
    return any(filename.lower().endswith(ext) for ext in ALLOWED_TEXT_EXTENSIONS)


# ---------------------------------------------------------------------------
# File conversion / extraction
# ---------------------------------------------------------------------------

def _image_to_pdf(image_bytes: bytes) -> bytes:
    from PIL import Image
    img = Image.open(io.BytesIO(image_bytes))
    if img.mode != "RGB":
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PDF", resolution=200)
    return buf.getvalue()


def _extract_image_base64(image_bytes: bytes) -> str:
    import base64
    return base64.b64encode(image_bytes).decode("ascii")


def _detect_image_media_type(image_bytes: bytes) -> str:
    if image_bytes[:8] == b'\x89PNG\r\n\x1a\n':
        return "image/png"
    if image_bytes[:2] == b'\xff\xd8':
        return "image/jpeg"
    if image_bytes[:4] in (b'II*\x00', b'MM\x00*'):
        return "image/tiff"
    if image_bytes[:4] == b'RIFF' and image_bytes[8:12] == b'WEBP':
        return "image/webp"
    if image_bytes[:2] == b'BM':
        return "image/bmp"
    return "image/png"


def _extract_text_file(content: bytes, filename: str) -> list[dict]:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("latin-1")

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    label = f"[{ext.upper()} FILE: {filename}]"

    if not text.strip():
        return [{"page": 1, "text": f"{label}\n(Empty file)"}]

    page_size = 3000
    pages = []
    for i in range(0, len(text), page_size):
        chunk = text[i:i + page_size]
        header = f"{label} (part {len(pages) + 1})" if len(text) > page_size else label
        pages.append({"page": len(pages) + 1, "text": f"{header}\n{chunk}"})

    return pages


def _extract_word_text(content: bytes) -> list[dict]:
    from docx import Document

    try:
        doc = Document(io.BytesIO(content))
    except Exception:
        raise HTTPException(status_code=400,
                            detail="Could not read Word document. Only .docx format is supported (not legacy .doc).")
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if "Heading" in style:
            parts.append(f"[HEADING] {text}")
        elif style.startswith("List"):
            parts.append(f"[LIST] {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")

    full_text = "\n".join(parts)
    if not full_text:
        return [{"page": 1, "text": "(Empty document)"}]

    page_size = 3000
    pages = []
    for i in range(0, len(full_text), page_size):
        chunk = full_text[i:i + page_size]
        pages.append({"page": len(pages) + 1, "text": chunk})

    return pages


def _extract_presentation_text(content: bytes) -> list[dict]:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(content))
    except Exception:
        raise HTTPException(status_code=400, detail="Could not read PowerPoint file. Only .pptx format is supported.")

    pages = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    rows.append(row_text)
                if rows:
                    texts.append("[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")
        pages.append({"page": i, "text": "\n".join(texts) if texts else "(Empty slide)"})

    return pages or [{"page": 1, "text": "(Empty presentation)"}]


def _extract_spreadsheet_text(content: bytes, filename: str) -> list[dict]:
    import pandas as pd

    pages = []
    lower = filename.lower()

    if lower.endswith(".csv"):
        df = pd.read_csv(io.BytesIO(content))
        text = df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
        pages.append({"page": 1, "text": f"[Sheet: CSV Data]\n\n{text}"})
    else:
        xls = pd.ExcelFile(io.BytesIO(content))
        for i, sheet_name in enumerate(xls.sheet_names, 1):
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if df.empty:
                continue
            text = df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
            pages.append({"page": i, "text": f"[Sheet: {sheet_name}]\n\n{text}"})

    if not pages:
        pages.append({"page": 1, "text": "(Empty spreadsheet)"})

    return pages


# ---------------------------------------------------------------------------
# Encryption helpers
# ---------------------------------------------------------------------------

def _encrypt_and_save(data: bytes, path: Path) -> None:
    if ENCRYPTION_KEY:
        from encryption import encrypt_bytes
        path.write_bytes(encrypt_bytes(data))
    else:
        path.write_bytes(data)


def _encrypt_text(text: str) -> str:
    if ENCRYPTION_KEY:
        from encryption import encrypt_bytes
        import base64
        return base64.b64encode(encrypt_bytes(text.encode("utf-8"))).decode("ascii")
    return text


def _decrypt_text(stored: str) -> str:
    if ENCRYPTION_KEY:
        from encryption import decrypt_bytes
        import base64
        return decrypt_bytes(base64.b64decode(stored)).decode("utf-8")
    return stored


def _safe_decrypt(stored: str, fallback: str = "") -> str:
    try:
        return _decrypt_text(stored)
    except Exception:
        return fallback


# ---------------------------------------------------------------------------
# PDF processing
# ---------------------------------------------------------------------------

def extract_text_from_bytes(pdf_bytes: bytes) -> list[dict]:
    return extract_text_with_ocr_fallback(pdf_bytes)


def _decrypt_and_load(filepath: Path) -> bytes:
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    if ENCRYPTION_KEY:
        from encryption import decrypt_file
        return decrypt_file(str(filepath))
    return filepath.read_bytes()


def _load_pdf_bytes(doc) -> bytes:
    filepath = Path(doc.filepath)
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="PDF file not found on disk")
    return _decrypt_and_load(filepath)


def _load_stored_text(doc) -> list[dict]:
    return json.loads(_safe_decrypt(doc.text_content, "[]"))


def _stored_text_to_structured(pages: list[dict]) -> str:
    parts = []
    for page_data in pages:
        parts.append(f"\n--- Page {page_data['page']} ---")
        if page_data.get("blocks"):
            for block in page_data["blocks"]:
                prefix = f"[{block['type'].upper()}] " if block['type'] != 'paragraph' else ""
                parts.append(f"{prefix}{block['text']}")
        elif page_data.get("text"):
            parts.append(page_data["text"])
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Upload validation
# ---------------------------------------------------------------------------

def validate_upload(file: UploadFile, content: bytes) -> tuple[str, bytes, str]:
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename is required")

    clean_name = _sanitize_filename(file.filename)
    is_image = _is_image_file(clean_name)
    is_spreadsheet = _is_spreadsheet_file(clean_name)
    is_word = _is_word_file(clean_name)
    is_presentation = _is_presentation_file(clean_name)
    is_text = _is_text_file(clean_name)

    if not clean_name.lower().endswith(".pdf") and not is_image and not is_spreadsheet and not is_word and not is_presentation and not is_text:
        raise HTTPException(status_code=400,
                            detail=f"Unsupported file type: {clean_name}")

    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=400,
            detail=f"{clean_name} exceeds max size of {MAX_FILE_SIZE // (1024*1024)} MB",
        )

    if is_spreadsheet:
        return clean_name, content, "spreadsheet"
    elif is_word:
        return clean_name, content, "word"
    elif is_presentation:
        return clean_name, content, "presentation"
    elif is_text:
        return clean_name, content, "text"
    elif is_image:
        try:
            pdf_bytes = _image_to_pdf(content)
        except Exception as exc:
            raise HTTPException(status_code=400,
                                detail=f"Could not process image {clean_name}: {exc}")
        return clean_name, pdf_bytes, "image"
    else:
        if not _verify_pdf_content(content):
            raise HTTPException(status_code=400,
                                detail="File does not appear to be a valid PDF")
        return clean_name, content, "pdf"


# ---------------------------------------------------------------------------
# Cleanup
# ---------------------------------------------------------------------------

def _run_cleanup(db) -> int:
    if not AUTO_CLEANUP_HOURS:
        return 0
    cutoff = datetime.now(timezone.utc) - timedelta(hours=AUTO_CLEANUP_HOURS)
    old_docs = db.query(DBDocument).filter(DBDocument.uploaded_at < cutoff).all()
    count = 0
    for doc in old_docs:
        _safe_unlink(Path(doc.filepath))
        db.delete(doc)
        count += 1
    if count:
        db.commit()
    return count


def _run_cleanup_background():
    db = SessionLocal()
    try:
        _run_cleanup(db)
    finally:
        db.close()


# ---------------------------------------------------------------------------
# Agent helpers (SSE formatting + Claude calls + cache)
# ---------------------------------------------------------------------------

def _agent_step(step: int, total: int, name: str, status: str = "running"):
    return f"data: {json.dumps({'type': 'agent_step', 'step': step, 'total': total, 'name': name, 'status': status})}\n\n"


def _agent_chunk(text: str):
    return f"data: {json.dumps({'type': 'chunk', 'text': text})}\n\n"


def _agent_done():
    return f"data: {json.dumps({'type': 'done'})}\n\n"


def _agent_error(msg: str):
    return f"data: {json.dumps({'type': 'error', 'detail': msg})}\n\n"


def _resolve_agent_model(model_choice: str) -> str:
    if model_choice in AGENT_MODELS:
        return AGENT_MODELS[model_choice]
    return CHAT_MODEL


def _call_claude(client, system: str, user_msg: str, tools=None, max_tokens=None, model=None):
    kwargs = dict(
        model=model or CHAT_MODEL,
        max_tokens=max_tokens or AGENT_MAX_TOKENS,
        system=system,
        messages=[{"role": "user", "content": user_msg}],
    )
    if tools:
        kwargs["tools"] = tools
    result = []
    with client.messages.stream(**kwargs) as stream:
        for event in stream:
            if hasattr(event, 'type') and event.type == 'content_block_delta':
                if hasattr(event.delta, 'text'):
                    result.append(event.delta.text)
    return "".join(result)


async def _call_claude_bg(client, system, user_msg, tools=None, max_tokens=None, model=None):
    import asyncio
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, lambda: _call_claude(client, system, user_msg, tools, max_tokens, model)
    )


def _stream_claude(client, system: str, user_msg: str, tools=None, max_tokens=None, model=None):
    kwargs = dict(
        model=model or CHAT_MODEL,
        max_tokens=max_tokens or AGENT_MAX_TOKENS,
        system=system,
        messages=[{"role": "user", "content": user_msg}],
    )
    if tools:
        kwargs["tools"] = tools
    with client.messages.stream(**kwargs) as stream:
        for event in stream:
            if hasattr(event, 'type') and event.type == 'content_block_delta':
                if hasattr(event.delta, 'text'):
                    yield event.delta.text


def _agent_cache_key(agent_type: str, model: str, doc_hashes: list[str], params: str) -> str:
    raw = f"{agent_type}|{model}|{'|'.join(sorted(doc_hashes))}|{params}"
    return hashlib.sha256(raw.encode()).hexdigest()


def _get_doc_hash(doc) -> str:
    if doc.content_hash:
        return doc.content_hash
    pdf_bytes = _load_pdf_bytes(doc)
    return hashlib.sha256(pdf_bytes).hexdigest()


def _check_agent_cache(db, cache_key: str, user_id: str | None = None):
    q = db.query(DBAgentCache).filter(DBAgentCache.cache_key == cache_key)
    if user_id:
        q = q.filter(DBAgentCache.user_id == user_id)
    else:
        q = q.filter(DBAgentCache.user_id.is_(None))
    cached = q.first()
    if not cached:
        return None
    if cached.expires_at and cached.expires_at < datetime.now(timezone.utc):
        db.delete(cached)
        db.commit()
        return None
    return _safe_decrypt(cached.result_data)


def _save_agent_cache(db, cache_key: str, agent_type: str, model: str,
                      result: str, doc_ids: list[str], params_summary: str,
                      user_id: str | None = None):
    docs = db.query(DBDocument).filter(DBDocument.id.in_(doc_ids)).all() if doc_ids else []
    cache = DBAgentCache(
        id=str(uuid.uuid4()),
        user_id=user_id,
        cache_key=cache_key,
        agent_type=agent_type,
        model_used=model,
        result_data=_encrypt_text(result),
        documents=docs,
        params_summary=params_summary[:200] if params_summary else "",
        created_at=datetime.now(timezone.utc),
        expires_at=datetime.now(timezone.utc) + timedelta(days=7),
    )
    db.add(cache)
    db.commit()


# ---------------------------------------------------------------------------
# Word document generation (shared by chat + updater routers)
# ---------------------------------------------------------------------------

def _markdown_to_docx(text: str, title: str = "Document") -> bytes:
    from docx import Document as DocxDocument
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = DocxDocument()

    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)

    t = doc.add_heading(title, level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = date_para.add_run(f"Generated: {datetime.now(timezone.utc).strftime('%B %d, %Y at %H:%M UTC')}")
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(128, 128, 128)
    doc.add_paragraph()

    lines = text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif re.match(r"^\d+[\.\)]\s", stripped):
            text_content = re.sub(r"^\d+[\.\)]\s", "", stripped)
            doc.add_paragraph(text_content, style="List Number")
        elif stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            p = doc.add_paragraph()
            run = p.add_run(stripped[2:-2])
            run.bold = True
            run.font.size = Pt(12)
        elif stripped in ("---", "***", "___"):
            doc.add_paragraph("_" * 50)
        elif not stripped:
            pass
        else:
            p = doc.add_paragraph()
            parts = re.split(r"(\*\*.*?\*\*|\*.*?\*)", stripped)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                    run = p.add_run(part[1:-1])
                    run.italic = True
                else:
                    p.add_run(part)
        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_vba_module(title: str) -> str:
    safe_title = title.replace('"', '""')
    return (
        'Attribute VB_Name = "ProcedureFormatter"\n'
        "' =============================================================================\n"
        "' Procedure Formatter VBA Module - Generated by PDFHelper AI\n"
        "'\n"
        "' HOW TO USE:\n"
        "' 1. Open the .docx file in Word\n"
        "' 2. Press Alt+F11 to open the VBA editor\n"
        "' 3. Go to Insert > Module, paste this code (or File > Import this .bas)\n"
        "' 4. Close VBA editor, press Alt+F8, run FormatProcedure\n"
        "' =============================================================================\n"
        "\n"
        "Public Sub FormatProcedure()\n"
        "    Application.ScreenUpdating = False\n"
        "    FormatHeadings\n"
        "    FormatStepNumbers\n"
        "    AddCheckboxes\n"
        "    FormatWarningsAndCautions\n"
        "    FormatTables\n"
        "    AddHeaderFooter\n"
        "    SetDocumentProperties\n"
        "    Application.ScreenUpdating = True\n"
        '    MsgBox "Procedure formatting complete!", vbInformation, "Procedure Formatter"\n'
        "End Sub\n"
        "\n"
        "Private Sub FormatHeadings()\n"
        "    Dim para As Paragraph\n"
        "    For Each para In ActiveDocument.Paragraphs\n"
        "        Dim txt As String\n"
        "        txt = Trim(para.Range.Text)\n"
        '        If UCase(txt) = txt And Len(txt) > 3 And Len(txt) < 80 And Not IsNumeric(Left(txt, 1)) Then\n'
        '            para.Style = ActiveDocument.Styles("Heading 1")\n'
        "        End If\n"
        '        If txt Like "#.0*" Or txt Like "##.0*" Or UCase(Left(txt, 7)) = "SECTION" Then\n'
        '            para.Style = ActiveDocument.Styles("Heading 1")\n'
        "        End If\n"
        '        If txt Like "#.#*" And Not txt Like "#.0*" Then\n'
        '            para.Style = ActiveDocument.Styles("Heading 2")\n'
        "        End If\n"
        '        If txt Like "#.#.#*" Then\n'
        '            para.Style = ActiveDocument.Styles("Heading 3")\n'
        "        End If\n"
        "    Next para\n"
        "End Sub\n"
        "\n"
        "Private Sub FormatStepNumbers()\n"
        "    Dim para As Paragraph\n"
        "    For Each para In ActiveDocument.Paragraphs\n"
        "        Dim txt As String\n"
        "        txt = Trim(para.Range.Text)\n"
        '        If UCase(Left(txt, 4)) = "STEP" Or txt Like "Step #*" Or txt Like "Step ##*" Then\n'
        "            para.Range.Font.Bold = True\n"
        "            para.Range.Font.Size = 11\n"
        "            para.Range.Font.Color = RGB(0, 51, 102)\n"
        "        End If\n"
        "    Next para\n"
        "End Sub\n"
        "\n"
        "Private Sub AddCheckboxes()\n"
        "    Dim para As Paragraph\n"
        "    For Each para In ActiveDocument.Paragraphs\n"
        "        Dim txt As String\n"
        "        txt = Trim(para.Range.Text)\n"
        '        If Left(txt, 3) = "[ ]" Then\n'
        '            para.Range.Text = ChrW(&H2610) & " " & Mid(txt, 4)\n'
        '        ElseIf UCase(Left(txt, 6)) = "VERIFY" Or UCase(Left(txt, 7)) = "CONFIRM" _\n'
        '               Or UCase(Left(txt, 5)) = "CHECK" Or UCase(Left(txt, 6)) = "ENSURE" Then\n'
        '            para.Range.InsertBefore ChrW(&H2610) & " "\n'
        "        End If\n"
        "    Next para\n"
        "End Sub\n"
        "\n"
        "Private Sub FormatWarningsAndCautions()\n"
        "    Dim para As Paragraph\n"
        "    For Each para In ActiveDocument.Paragraphs\n"
        "        Dim txt As String\n"
        "        txt = UCase(Trim(para.Range.Text))\n"
        '        If Left(txt, 7) = "WARNING" Or Left(txt, 6) = "DANGER" Then\n'
        "            para.Range.Font.Bold = True\n"
        "            para.Range.Font.Color = RGB(204, 0, 0)\n"
        "            para.Shading.BackgroundPatternColor = RGB(255, 235, 235)\n"
        '        ElseIf Left(txt, 7) = "CAUTION" Then\n'
        "            para.Range.Font.Bold = True\n"
        "            para.Range.Font.Color = RGB(204, 102, 0)\n"
        "            para.Shading.BackgroundPatternColor = RGB(255, 248, 230)\n"
        '        ElseIf Left(txt, 4) = "NOTE" Then\n'
        "            para.Range.Font.Italic = True\n"
        "            para.Range.Font.Color = RGB(0, 51, 153)\n"
        "            para.Shading.BackgroundPatternColor = RGB(235, 243, 255)\n"
        "        End If\n"
        "    Next para\n"
        "End Sub\n"
        "\n"
        "Private Sub FormatTables()\n"
        "    Dim tbl As Table\n"
        "    For Each tbl In ActiveDocument.Tables\n"
        "        tbl.Borders.Enable = True\n"
        "        tbl.Borders.InsideLineStyle = wdLineStyleSingle\n"
        "        tbl.Borders.OutsideLineStyle = wdLineStyleSingle\n"
        "        If tbl.Rows.Count > 0 Then\n"
        "            tbl.Rows(1).Range.Font.Bold = True\n"
        "            tbl.Rows(1).Shading.BackgroundPatternColor = RGB(0, 51, 102)\n"
        "            tbl.Rows(1).Range.Font.Color = RGB(255, 255, 255)\n"
        "        End If\n"
        "        Dim i As Long\n"
        "        For i = 2 To tbl.Rows.Count\n"
        "            If i Mod 2 = 0 Then\n"
        "                tbl.Rows(i).Shading.BackgroundPatternColor = RGB(242, 246, 250)\n"
        "            End If\n"
        "        Next i\n"
        "        tbl.AutoFitBehavior wdAutoFitWindow\n"
        "    Next tbl\n"
        "End Sub\n"
        "\n"
        "Private Sub AddHeaderFooter()\n"
        "    Dim sec As Section\n"
        "    For Each sec In ActiveDocument.Sections\n"
        "        sec.Headers(wdHeaderFooterPrimary).Range.Text = _\n"
        f'            "{safe_title}" & vbTab & vbTab & "CONTROLLED DOCUMENT"\n'
        "        sec.Headers(wdHeaderFooterPrimary).Range.Font.Size = 9\n"
        "        sec.Headers(wdHeaderFooterPrimary).Range.Font.Color = RGB(128, 128, 128)\n"
        "        sec.Footers(wdHeaderFooterPrimary).Range.Text = _\n"
        '            "Page " & vbTab & vbTab & "Revision Date: " & Format(Date, "yyyy-mm-dd")\n'
        "        sec.Footers(wdHeaderFooterPrimary).Range.Font.Size = 9\n"
        "        sec.Footers(wdHeaderFooterPrimary).Range.Font.Color = RGB(128, 128, 128)\n"
        "        Dim rng As Range\n"
        "        Set rng = sec.Footers(wdHeaderFooterPrimary).Range\n"
        "        rng.Collapse Direction:=wdCollapseStart\n"
        "        rng.MoveEnd Unit:=wdCharacter, Count:=5\n"
        "        rng.Collapse Direction:=wdCollapseEnd\n"
        "        ActiveDocument.Fields.Add Range:=rng, Type:=wdFieldPage\n"
        '        rng.InsertAfter " of "\n'
        "        rng.Collapse Direction:=wdCollapseEnd\n"
        "        ActiveDocument.Fields.Add Range:=rng, Type:=wdFieldNumPages\n"
        "    Next sec\n"
        "End Sub\n"
        "\n"
        "Private Sub SetDocumentProperties()\n"
        "    With ActiveDocument.BuiltInDocumentProperties\n"
        f'        .Item("Title").Value = "{safe_title}"\n'
        '        .Item("Subject").Value = "Operating Procedure"\n'
        '        .Item("Category").Value = "Procedure Document"\n'
        "    End With\n"
        '    ActiveDocument.Styles("Normal").Font.Name = "Calibri"\n'
        '    ActiveDocument.Styles("Normal").Font.Size = 11\n'
        "End Sub\n"
        "\n"
        "Public Sub InsertRevisionTable()\n"
        "    Dim tbl As Table\n"
        "    Set tbl = ActiveDocument.Tables.Add( _\n"
        "        Range:=Selection.Range, NumRows:=4, NumColumns:=4)\n"
        '    tbl.Cell(1, 1).Range.Text = "Rev"\n'
        '    tbl.Cell(1, 2).Range.Text = "Date"\n'
        '    tbl.Cell(1, 3).Range.Text = "Description"\n'
        '    tbl.Cell(1, 4).Range.Text = "Author"\n'
        '    tbl.Cell(2, 1).Range.Text = "0"\n'
        '    tbl.Cell(2, 2).Range.Text = Format(Date, "yyyy-mm-dd")\n'
        '    tbl.Cell(2, 3).Range.Text = "Initial release - AI-generated from source documents"\n'
        '    tbl.Cell(2, 4).Range.Text = ""\n'
        "    tbl.Rows(1).Range.Font.Bold = True\n"
        "    tbl.Rows(1).Shading.BackgroundPatternColor = RGB(0, 51, 102)\n"
        "    tbl.Rows(1).Range.Font.Color = RGB(255, 255, 255)\n"
        "    tbl.Borders.Enable = True\n"
        "    tbl.AutoFitBehavior wdAutoFitWindow\n"
        "End Sub\n"
        "\n"
        "Public Sub InsertSignOffBlock()\n"
        "    Dim tbl As Table\n"
        "    Set tbl = ActiveDocument.Tables.Add( _\n"
        "        Range:=Selection.Range, NumRows:=4, NumColumns:=3)\n"
        '    tbl.Cell(1, 1).Range.Text = "Role"\n'
        '    tbl.Cell(1, 2).Range.Text = "Name / Signature"\n'
        '    tbl.Cell(1, 3).Range.Text = "Date"\n'
        '    tbl.Cell(2, 1).Range.Text = "Prepared By"\n'
        '    tbl.Cell(3, 1).Range.Text = "Reviewed By"\n'
        '    tbl.Cell(4, 1).Range.Text = "Approved By"\n'
        "    tbl.Rows(1).Range.Font.Bold = True\n"
        "    tbl.Rows(1).Shading.BackgroundPatternColor = RGB(0, 51, 102)\n"
        "    tbl.Rows(1).Range.Font.Color = RGB(255, 255, 255)\n"
        "    tbl.Borders.Enable = True\n"
        "    Dim i As Long\n"
        "    For i = 2 To 4\n"
        "        tbl.Rows(i).Height = CentimetersToPoints(1.5)\n"
        "    Next i\n"
        "    tbl.AutoFitBehavior wdAutoFitWindow\n"
        "End Sub\n"
    )
